from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
THEMES_DIR = ROOT / "themes"
PHILIPS_ASSETS = ROOT.parent.parent / "philips" / "philips-assets"
OUTPUT_PPTX = ROOT / "philips-theme-showcase.pptx"

WORDMARK_PATH = PHILIPS_ASSETS / "Logos" / "philips_wordmark" / "GMC_Wordmark_2008_RGB.jpg"
SHIELD_RGB_PATH = PHILIPS_ASSETS / "Logos" / "philips-shield" / "philips shield_PNG" / "Shield_RGB_2014.png"
SHIELD_WHITE_PATH = PHILIPS_ASSETS / "Logos" / "philips-shield" / "philips shield_PNG" / "Shield_White_2014.png"
WEBINAR_LIGHT_PATH = PHILIPS_ASSETS / "Webinar backgrounds" / "webinar_background_light.png"
WEBINAR_DARK_PATH = PHILIPS_ASSETS / "Webinar backgrounds" / "webinar_background_dark.png"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

THEME_CONFIG = {
    "philips-powerpoint": {
        "page_bg": "#F4F8FC",
        "ink": "#00126E",
        "preview": "powerpoint",
    },
    "philips-webinar-light": {
        "page_bg": "#EAF6FD",
        "ink": "#04548E",
        "preview": WEBINAR_LIGHT_PATH,
        "shield": SHIELD_WHITE_PATH,
    },
    "philips-webinar-dark": {
        "page_bg": "#0D4A85",
        "ink": "#FFFFFF",
        "preview": WEBINAR_DARK_PATH,
        "shield": SHIELD_WHITE_PATH,
    },
}


def hex_to_rgb(hex_color: str) -> RGBColor:
    value = hex_color.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_textbox(slide, left, top, width, height, text, font_name, font_size, color, bold=False, align=PP_ALIGN.LEFT):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return textbox


def set_slide_background(slide, hex_color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(hex_color)


def parse_theme(theme_path: Path) -> dict[str, object]:
    text = theme_path.read_text(encoding="utf-8")
    lines = text.splitlines()
    title = lines[0].lstrip("# ").strip()

    sections: dict[str, list[str]] = {"intro": []}
    current = "intro"
    for line in lines[1:]:
        if line.startswith("## "):
            current = line[3:].strip().lower()
            sections[current] = []
        else:
            sections.setdefault(current, []).append(line)

    description = " ".join(line.strip() for line in sections.get("intro", []) if line.strip())
    palette = []
    for line in sections.get("color palette", []):
        if not line.startswith("- "):
            continue
        match = re.match(r"- \*\*(.+?)\*\*: `?(#?[0-9A-Fa-f]{6})`? - (.+)", line.strip())
        if not match:
            continue
        name, color_hex, note = match.groups()
        if not color_hex.startswith("#"):
            color_hex = f"#{color_hex}"
        palette.append({"name": name, "hex": color_hex.upper(), "note": note})

    typography = []
    for line in sections.get("typography", []):
        if line.startswith("- ") and ": " in line:
            label, value = line[2:].split(": ", 1)
            typography.append((label.replace("**", ""), value))

    best_used_for = " ".join(line.strip() for line in sections.get("best used for", []) if line.strip())
    return {
        "id": theme_path.stem,
        "title": title,
        "description": description,
        "palette": palette,
        "typography": typography,
        "best_used_for": best_used_for,
    }


def add_round_rect(slide, left, top, width, height, fill_hex, line_hex=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = hex_to_rgb(line_hex)
    else:
        shape.line.fill.background()
    return shape


def add_cover_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, "#0B5ED7")
    if WORDMARK_PATH.exists():
        slide.shapes.add_picture(str(WORDMARK_PATH), Inches(10.7), Inches(0.42), width=Inches(2.1), height=Inches(0.38))

    add_textbox(slide, Inches(0.7), Inches(0.62), Inches(6.8), Inches(0.6), "Philips Theme Showcase", "Neue Frutiger World Thin", 30, RGBColor(255, 255, 255))
    add_textbox(slide, Inches(0.72), Inches(1.1), Inches(7.5), Inches(0.6), "Theme Factory companion for Philips-branded presentations and digital artifacts", "Neue Frutiger World Book", 14, RGBColor(255, 255, 255))

    card_top = Inches(2.0)
    card_height = Inches(4.65)
    card_width = Inches(3.75)
    card_gap = Inches(0.28)
    cards = [
        ("Philips PowerPoint", "#FFFFFF", "#0B5ED7", "#00126E"),
        ("Webinar Light", "#EAF6FD", "#1780BD", "#04548E"),
        ("Webinar Dark", "#0F4276", "#1490CF", "#FFFFFF"),
    ]
    for index, (label, bg_hex, accent_hex, text_hex) in enumerate(cards):
        left = Inches(0.72) + index * (card_width + card_gap)
        add_round_rect(slide, left, card_top, card_width, card_height, bg_hex)
        add_round_rect(slide, left + Inches(0.18), card_top + Inches(0.22), card_width - Inches(0.36), Inches(0.72), accent_hex)
        add_textbox(slide, left + Inches(0.34), card_top + Inches(0.4), card_width - Inches(0.68), Inches(0.3), label, "Neue Frutiger World Light", 19, hex_to_rgb(text_hex))
        add_round_rect(slide, left + Inches(0.18), card_top + Inches(1.2), card_width - Inches(0.36), Inches(3.05), "#F7FAFD" if index == 0 else "#FFFFFF")
        if index == 0:
            add_round_rect(slide, left + Inches(0.38), card_top + Inches(1.46), Inches(1.05), Inches(1.8), "#BDF0FF")
            add_round_rect(slide, left + Inches(1.58), card_top + Inches(1.46), Inches(0.52), Inches(1.8), "#00126E")
            add_round_rect(slide, left + Inches(2.23), card_top + Inches(1.46), Inches(0.82), Inches(1.8), "#DCE6F5")
        else:
            add_round_rect(slide, left + Inches(0.38), card_top + Inches(1.5), Inches(2.7), Inches(1.55), accent_hex)


def add_theme_slide(prs: Presentation, theme: dict[str, object]) -> None:
    config = THEME_CONFIG[theme["id"]]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, config["page_bg"])

    if WORDMARK_PATH.exists():
        slide.shapes.add_picture(str(WORDMARK_PATH), Inches(0.55), Inches(0.35), width=Inches(1.7), height=Inches(0.3))

    add_textbox(slide, Inches(0.62), Inches(0.82), Inches(4.8), Inches(0.5), theme["title"], "Neue Frutiger World Thin", 27, hex_to_rgb(config["ink"]))
    add_textbox(slide, Inches(0.64), Inches(1.26), Inches(4.6), Inches(0.8), theme["description"], "Neue Frutiger World Light", 13, hex_to_rgb(config["ink"]))

    preview_left = Inches(6.65)
    preview_top = Inches(0.72)
    preview_width = Inches(6.0)
    preview_height = Inches(6.0)

    if config["preview"] == "powerpoint":
        add_round_rect(slide, preview_left, preview_top, preview_width, preview_height, "#FFFFFF")
        add_round_rect(slide, preview_left + Inches(0.2), preview_top + Inches(0.22), preview_width - Inches(0.4), Inches(0.76), "#0B5ED7")
        add_textbox(slide, preview_left + Inches(0.42), preview_top + Inches(0.43), Inches(4.4), Inches(0.3), "Philips PowerPoint", "Neue Frutiger World Thin", 23, RGBColor(255, 255, 255))
        add_textbox(slide, preview_left + Inches(0.38), preview_top + Inches(1.18), Inches(4.8), Inches(0.45), "Brand blue titles with spacious, restrained content framing", "Neue Frutiger World Light", 15, hex_to_rgb("#00126E"))
        add_round_rect(slide, preview_left + Inches(0.38), preview_top + Inches(1.8), Inches(2.0), Inches(2.9), "#BDF0FF")
        add_round_rect(slide, preview_left + Inches(2.56), preview_top + Inches(1.8), Inches(0.88), Inches(2.9), "#00126E")
        add_round_rect(slide, preview_left + Inches(3.62), preview_top + Inches(1.8), Inches(1.38), Inches(2.9), "#DCE6F5")
        if SHIELD_RGB_PATH.exists():
            slide.shapes.add_picture(str(SHIELD_RGB_PATH), preview_left + Inches(4.94), preview_top + Inches(0.26), width=Inches(0.62), height=Inches(0.84))
    else:
        if Path(config["preview"]).exists():
            slide.shapes.add_picture(str(config["preview"]), preview_left, preview_top, width=preview_width, height=preview_height)
        border = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, preview_left, preview_top, preview_width, preview_height)
        border.fill.background()
        border.line.color.rgb = RGBColor(255, 255, 255)
        if Path(config["shield"]).exists():
            slide.shapes.add_picture(str(config["shield"]), preview_left + Inches(5.1), preview_top + Inches(0.24), width=Inches(0.6), height=Inches(0.82))

    section_color = hex_to_rgb("#5D6B7A") if config["ink"] != "#FFFFFF" else hex_to_rgb("#D7E8F7")
    add_textbox(slide, Inches(0.66), Inches(2.15), Inches(2.0), Inches(0.2), "COLOR PALETTE", "Neue Frutiger World Book", 9, section_color)

    swatch_top = Inches(2.45)
    for index, color_spec in enumerate(theme["palette"]):
        y = swatch_top + Inches(index * 0.68)
        swatch = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.68), y, Inches(0.22), Inches(0.22))
        swatch.fill.solid()
        swatch.fill.fore_color.rgb = hex_to_rgb(color_spec["hex"])
        swatch.line.fill.background()
        add_textbox(slide, Inches(1.0), y - Inches(0.03), Inches(3.9), Inches(0.2), f"{color_spec['name']} {color_spec['hex']}", "Neue Frutiger World Light", 11, hex_to_rgb(config["ink"]))
        add_textbox(slide, Inches(1.0), y + Inches(0.17), Inches(4.5), Inches(0.38), color_spec["note"], "Neue Frutiger World Book", 9, hex_to_rgb(config["ink"]))

    add_textbox(slide, Inches(0.66), Inches(5.15), Inches(2.0), Inches(0.2), "TYPOGRAPHY", "Neue Frutiger World Book", 9, section_color)
    typography_lines = "\n".join(f"{label}: {value}" for label, value in theme["typography"])
    typography_box = slide.shapes.add_textbox(Inches(0.66), Inches(5.42), Inches(4.7), Inches(0.8))
    frame = typography_box.text_frame
    frame.word_wrap = True
    for idx, line in enumerate(typography_lines.splitlines()):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.font.name = "Neue Frutiger World Book"
        paragraph.font.size = Pt(11)
        paragraph.font.color.rgb = hex_to_rgb(config["ink"])

    add_textbox(slide, Inches(0.66), Inches(6.2), Inches(2.2), Inches(0.2), "BEST USED FOR", "Neue Frutiger World Book", 9, section_color)
    add_textbox(slide, Inches(0.66), Inches(6.46), Inches(5.2), Inches(0.48), theme["best_used_for"], "Neue Frutiger World Book", 10, hex_to_rgb(config["ink"]))


def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    add_cover_slide(prs)
    for theme_name in ["philips-powerpoint", "philips-webinar-light", "philips-webinar-dark"]:
        add_theme_slide(prs, parse_theme(THEMES_DIR / f"{theme_name}.md"))
    prs.save(str(OUTPUT_PPTX))


if __name__ == "__main__":
    main()